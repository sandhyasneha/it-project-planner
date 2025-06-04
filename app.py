import os
import streamlit as st
import sqlite3
import hashlib
from openai import OpenAI
import pyttsx3
import pyperclip
import speech_recognition as sr
from datetime import datetime
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from pptx import Presentation
from pptx.util import Inches

# ----- CONFIG -----
st.set_page_config(page_title="IT Project Planner", page_icon="🛠️")
api_key = os.getenv("OPENAI_API_KEY")
st.write(f"API Key loaded: {api_key is not None}")
client = OpenAI(api_key=api_key)

# ----- DB HELPERS -----
def make_hashes(password):
    return hashlib.sha256(password.encode()).hexdigest()

def verify_hashes(password, hashed_text):
    return make_hashes(password) == hashed_text

def create_usertable():
    conn = sqlite3.connect('users.db')
    conn.execute('CREATE TABLE IF NOT EXISTS userstable(email TEXT, password TEXT)')
    conn.commit()
    conn.close()

def add_userdata(email, password):
    conn = sqlite3.connect('users.db')
    conn.execute('INSERT INTO userstable(email, password) VALUES (?, ?)', (email, make_hashes(password)))
    conn.commit()
    conn.close()

def login_user(email, password):
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute('SELECT * FROM userstable WHERE email = ? AND password = ?', (email, make_hashes(password)))
    data = c.fetchall()
    conn.close()
    return data

def get_all_users():
    conn = sqlite3.connect('users.db')
    c = conn.cursor()
    c.execute('SELECT email FROM userstable')
    data = c.fetchall()
    conn.close()
    return [email[0] for email in data]

# ----- EMAIL HELPER -----
def send_reminder_to_all():
    sender_email = os.getenv("SENDER_EMAIL")
    sender_password = os.getenv("SENDER_PASSWORD")
    users = get_all_users()

    msg = MIMEMultipart()
    msg['From'] = sender_email
    msg['Subject'] = "Weekly Timesheet Reminder"
    body = "Team, gentle reminder to update your time sheet without fail."
    msg.attach(MIMEText(body, 'plain'))

    try:
        with smtplib.SMTP_SSL('smtp.gmail.com', 465) as server:
            server.login(sender_email, sender_password)
            for email in users:
                msg['To'] = email
                server.sendmail(sender_email, email, msg.as_string())
        st.success("Reminder sent to all users.")
    except Exception as e:
        st.error(f"Error sending emails: {e}")

# ----- PPT GENERATOR -----
def generate_ppt(subject, content):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = subject
    body.text = content
    file_path = f"{subject.replace(' ', '_')}.pptx"
    prs.save(file_path)
    return file_path

# ----- LOGIN -----
create_usertable()
menu = ["Login", "SignUp"]
if 'logged_in' not in st.session_state:
    st.session_state.logged_in = False
    st.session_state.user_email = ""
choice = st.sidebar.selectbox("Menu", menu)

if choice == "Login":
    st.sidebar.subheader("Login")
    email = st.sidebar.text_input("Email")
    password = st.sidebar.text_input("Password", type="password")
    if st.sidebar.button("Login"):
        if not email.endswith("@nttdata.com"):
            st.sidebar.error("Only @nttdata.com emails allowed")
        elif login_user(email, password):
            st.session_state.logged_in = True
            st.session_state.user_email = email
            st.sidebar.success(f"Welcome {email}")
        else:
            st.sidebar.error("Invalid credentials")

elif choice == "SignUp":
    st.sidebar.subheader("Create Account")
    email = st.sidebar.text_input("New Email")
    password = st.sidebar.text_input("New Password", type="password")
    if st.sidebar.button("Create Account"):
        if not email.endswith("@nttdata.com"):
            st.sidebar.error("Only @nttdata.com emails allowed")
        else:
            add_userdata(email, password)
            st.sidebar.success("Account created!")

if st.session_state.logged_in:
    if st.sidebar.button("Logout"):
        st.session_state.logged_in = False
        st.session_state.user_email = ""
        st.experimental_rerun()

if st.session_state.logged_in:
    user_email = st.session_state.user_email
    st.title("🛠️ IT Project Planner")

    if user_email.lower() == "george@nttdata.com" and datetime.today().weekday() != 4:
        st.sidebar.info("Hi George, today is not Friday. Do you want to send a timesheet reminder?")
        if st.sidebar.button("Yes, Send Reminder"):
            st.sidebar.success("Reminder message sent.")
        elif st.sidebar.button("No"):
            st.sidebar.info("Cancelled.")

    if 'input_text' not in st.session_state:
        st.session_state.input_text = ""
    if 'generated_plan' not in st.session_state:
        st.session_state.generated_plan = ""

    st.session_state.input_text = st.text_area("Describe your project:", st.session_state.input_text)

    if st.button("Generate Plan") and st.session_state.input_text:
        with st.spinner("Generating plan..."):
            try:
                response = client.chat.completions.create(
                    model="gpt-3.5-turbo",
                    messages=[
                        {"role": "system", "content": "You are an expert IT project manager."},
                        {"role": "user", "content": st.session_state.input_text}
                    ],
                    max_tokens=800
                )
                st.session_state.generated_plan = response.choices[0].message.content
            except Exception as e:
                st.error(f"Error: {e}")

    if st.session_state.generated_plan:
        st.markdown(st.session_state.generated_plan)

        if st.button("📋 Copy Plan"):
            try:
                pyperclip.copy(st.session_state.generated_plan)
                st.success("Copied to clipboard!")
            except Exception as e:
                st.error("Clipboard error: Please install xclip or xsel on Linux.")

        if st.button("🔊 Play Plan"):
            engine = pyttsx3.init()
            engine.say(st.session_state.generated_plan)
            engine.runAndWait()

        st.download_button("📥 Download Plan", st.session_state.generated_plan, file_name="plan.txt")

        # Generate PPT
        if st.button("📊 Generate PowerPoint"):
            ppt_file = generate_ppt("Project Plan", st.session_state.generated_plan)
            with open(ppt_file, "rb") as f:
                st.download_button("Download PPT", f, file_name=ppt_file, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    if st.button("🎙️ Dictate (local mic only)"):
        try:
            recognizer = sr.Recognizer()
            with sr.Microphone() as source:
                st.info("Listening...")
                audio = recognizer.listen(source)
            text = recognizer.recognize_google(audio)
            st.session_state.input_text = text
            st.success(f"Recognized: {text}")
        except Exception as e:
            st.error(f"Voice error: {e}")
